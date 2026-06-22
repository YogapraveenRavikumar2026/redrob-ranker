import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def main():
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    dark_navy = RGBColor(15, 27, 45)      # #0F1B2D
    teal = RGBColor(0, 201, 167)          # #00C9A7
    white = RGBColor(255, 255, 255)
    light_gray = RGBColor(200, 208, 220)
    card_bg = RGBColor(22, 37, 59)         # #16253B
    card_border = RGBColor(38, 57, 84)     # #263954
    
    # Helper to create slide with solid background
    def create_slide():
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = dark_navy
        return slide

    # Helper to add standard slide header
    def add_header(slide, title_text, category_text=None):
        if category_text:
            cat_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.333), Inches(0.35))
            tf = cat_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = category_text.upper()
            p.font.name = 'Segoe UI'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = teal
            
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.65), Inches(11.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = white

    # Helper to format table cells
    def format_cell(cell, text, font_size=12, bold=False, text_color=white, bg_color=None, align=PP_ALIGN.LEFT):
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Segoe UI'
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color
        p.alignment = align
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # -------------------------------------------------------------
    # SLIDE 1 — Title
    # -------------------------------------------------------------
    slide1 = create_slide()
    # Accent colored left border line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.2), Inches(0.08), Inches(2.8))
    line.fill.solid()
    line.fill.fore_color.rgb = teal
    line.line.color.rgb = teal
    
    # Title & Subtitle box
    title_box = slide1.shapes.add_textbox(Inches(1.3), Inches(2.1), Inches(10.5), Inches(2.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "Recruiter Brain"
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = teal
    p1.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "AI Candidate Ranking That Thinks Like a Great Recruiter"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(22)
    p2.font.color.rgb = white
    
    # Bottom metadata box
    meta_box = slide1.shapes.add_textbox(Inches(1.3), Inches(5.1), Inches(10.5), Inches(0.6))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    tf_meta.margin_left = tf_meta.margin_top = tf_meta.margin_right = tf_meta.margin_bottom = 0
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Redrob Hackathon 2026 | Team Antigravity"
    p_meta.font.name = 'Segoe UI'
    p_meta.font.size = Pt(14)
    p_meta.font.bold = True
    p_meta.font.color.rgb = light_gray

    # -------------------------------------------------------------
    # SLIDE 2 — The Problem
    # -------------------------------------------------------------
    slide2 = create_slide()
    add_header(slide2, "Why keyword matching fails", "The Core Problem")
    
    # Text box for bullet points
    bullets_box = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.2))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    bullets = [
        "Keyword filters miss Tier-5 candidates who built RAG before it had a name",
        "A perfect skill list means nothing if the candidate hasn't logged in for 6 months",
        "HR Managers with 9 AI keywords rank above ML Engineers who shipped real systems"
    ]
    for idx, b_text in enumerate(bullets):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = "•   " + b_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(18)
        p.font.color.rgb = white
        p.space_after = Pt(24)
        
    # Callout box
    callout_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.1))
    callout_bg.fill.solid()
    callout_bg.fill.fore_color.rgb = card_bg
    callout_bg.line.color.rgb = teal
    callout_bg.line.width = Pt(1.5)
    
    tf_call = callout_bg.text_frame
    tf_call.word_wrap = True
    tf_call.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_call.margin_left = Inches(0.3)
    p_call = tf_call.paragraphs[0]
    p_call.text = "The sample submission proves this — rank 1 is an HR Manager"
    p_call.font.name = 'Segoe UI'
    p_call.font.size = Pt(16)
    p_call.font.bold = True
    p_call.font.color.rgb = teal

    # -------------------------------------------------------------
    # SLIDE 3 — Our Insight
    # -------------------------------------------------------------
    slide3 = create_slide()
    add_header(slide3, "What great recruiters actually do", "Our Insight")
    
    # Left Column Card (Keyword Matching)
    left_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(5.4), Inches(4.8))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = card_bg
    left_bg.line.color.rgb = card_border
    
    tf_l = left_bg.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.4)
    tf_l.margin_top = Inches(0.4)
    tf_l.margin_right = Inches(0.4)
    
    p_lt = tf_l.paragraphs[0]
    p_lt.text = "❌ Keyword Matching"
    p_lt.font.name = 'Segoe UI'
    p_lt.font.size = Pt(20)
    p_lt.font.bold = True
    p_lt.font.color.rgb = RGBColor(239, 83, 80) # Light red
    p_lt.space_after = Pt(20)
    
    l_points = [
        "Matches literal words, ignoring synonyms and context",
        "Treats all skills equally regardless of depth",
        "Ignores reachability and availability signals",
        "Easily gamed and fooled by keyword stuffers"
    ]
    for lp in l_points:
        p = tf_l.add_paragraph()
        p.text = "•   " + lp
        p.font.name = 'Segoe UI'
        p.font.size = Pt(14)
        p.font.color.rgb = light_gray
        p.space_after = Pt(16)
        
    # Right Column Card (Recruiter Brain)
    right_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.4), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = card_bg
    right_bg.line.color.rgb = teal
    right_bg.line.width = Pt(1.5)
    
    tf_r = right_bg.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.4)
    tf_r.margin_top = Inches(0.4)
    tf_r.margin_right = Inches(0.4)
    
    p_rt = tf_r.paragraphs[0]
    p_rt.text = "✅ Recruiter Brain"
    p_rt.font.name = 'Segoe UI'
    p_rt.font.size = Pt(20)
    p_rt.font.bold = True
    p_rt.font.color.rgb = teal
    p_rt.space_after = Pt(20)
    
    r_points = [
        "Reads career trajectory and progression patterns",
        "Weighs production experience vs pure research",
        "Checks availability, response rate, and notice days",
        "Detects and filters out impossible profiles and honeypots"
    ]
    for rp in r_points:
        p = tf_r.add_paragraph()
        p.text = "•   " + rp
        p.font.name = 'Segoe UI'
        p.font.size = Pt(14)
        p.font.color.rgb = white
        p.space_after = Pt(16)

    # -------------------------------------------------------------
    # SLIDE 4 — Architecture
    # -------------------------------------------------------------
    slide4 = create_slide()
    add_header(slide4, "5-layer pipeline", "Architecture")
    
    layers = [
        ("Data Ingestion", "100K candidates + JD loaded into RAM"),
        ("Feature Extraction", "40+ signals per candidate calculated in parallel on all 100K"),
        ("Pre-filter", "Hard rules eliminate 97% of non-fits in 28 seconds"),
        ("Deep Scoring", "Semantic embeddings + 5-axis scoring computed on top 3000 survivors"),
        ("Output", "Top 100 with unique, rank-consistent natural language reasoning")
    ]
    
    top_pos = 1.7
    for idx, (title, desc) in enumerate(layers, 1):
        # Circle with step number
        circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(top_pos), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = teal
        circle.line.fill.background()
        
        tf_c = circle.text_frame
        tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_c = tf_c.paragraphs[0]
        p_c.text = str(idx)
        p_c.font.name = 'Segoe UI'
        p_c.font.size = Pt(16)
        p_c.font.bold = True
        p_c.font.color.rgb = dark_navy
        p_c.alignment = PP_ALIGN.CENTER
        
        # Text block
        tb = slide4.shapes.add_textbox(Inches(1.8), Inches(top_pos - 0.05), Inches(10.5), Inches(0.7))
        tf_tb = tb.text_frame
        tf_tb.word_wrap = True
        tf_tb.margin_left = tf_tb.margin_top = tf_tb.margin_right = tf_tb.margin_bottom = 0
        
        p_title = tf_tb.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Segoe UI'
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = teal
        
        p_desc = tf_tb.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Segoe UI'
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = white
        
        top_pos += 1.05

    # -------------------------------------------------------------
    # SLIDE 5 — The Scoring Model
    # -------------------------------------------------------------
    slide5 = create_slide()
    add_header(slide5, "5-axis scoring — not just keywords", "Scoring Engine")
    
    table_shape = slide5.shapes.add_table(6, 3, Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(7.033)
    
    headers = ["Axis", "Weight", "What it measures"]
    for i, h_text in enumerate(headers):
        format_cell(table.cell(0, i), h_text, font_size=15, bold=True, text_color=dark_navy, bg_color=teal, align=PP_ALIGN.CENTER if i==1 else PP_ALIGN.LEFT)
        
    data = [
        ("Semantic fit", "35%", "sentence-transformers cosine similarity between JD & resume summary"),
        ("Hard requirements", "30%", "embeddings, vector DB, ranking eval, production ML, and python depth"),
        ("Career trajectory", "15%", "product company experience, recent role relevance weights (Last role = 3x)"),
        ("Availability", "15%", "active last 90d, open to work status, response rate, notice period"),
        ("Location fit", "5%", "Pune/Noida exact match vs India tier-1 vs India other vs outside")
    ]
    
    for row_idx, (axis, weight, measure) in enumerate(data, 1):
        cell_bg = card_bg if row_idx % 2 == 0 else dark_navy
        format_cell(table.cell(row_idx, 0), axis, font_size=13, bold=True, text_color=white, bg_color=cell_bg)
        format_cell(table.cell(row_idx, 1), weight, font_size=13, bold=True, text_color=teal, bg_color=cell_bg, align=PP_ALIGN.CENTER)
        format_cell(table.cell(row_idx, 2), measure, font_size=13, bold=False, text_color=light_gray, bg_color=cell_bg)

    # -------------------------------------------------------------
    # SLIDE 6 — Honeypot Defense
    # -------------------------------------------------------------
    slide6 = create_slide()
    add_header(slide6, "We catch what others miss", "Honeypot Defense")
    
    sub = slide6.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(0.4))
    tf_sub = sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "5 detection rules for impossible profiles:"
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = teal
    
    bullets_box = slide6.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.333), Inches(3.3))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    rules = [
        "Expert skill explosion: 8+ expert skills with <5 years experience",
        "Title-skill mismatch: Marketing Manager with 5+ AI expert skills",
        "Date inconsistency: role end_date before start_date",
        "Career duration overflow: simple duration sum > years_experience × 12 + 24",
        "Ghost profile: 95%+ complete, zero activity, no GitHub, 0 applications"
    ]
    for idx, rule in enumerate(rules):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = "•   " + rule
        p.font.name = 'Segoe UI'
        p.font.size = Pt(15)
        p.font.color.rgb = white
        p.space_after = Pt(14)
        
    callout = slide6.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.333), Inches(0.7))
    tf_call = callout.text_frame
    tf_call.word_wrap = True
    tf_call.margin_left = tf_call.margin_top = tf_call.margin_right = tf_call.margin_bottom = 0
    p_call = tf_call.paragraphs[0]
    p_call.text = "Hard honeypots → score multiplied by 0.05 (effectively eliminated)"
    p_call.font.name = 'Segoe UI'
    p_call.font.size = Pt(16)
    p_call.font.bold = True
    p_call.font.color.rgb = teal

    # -------------------------------------------------------------
    # SLIDE 7 — Performance
    # -------------------------------------------------------------
    slide7 = create_slide()
    add_header(slide7, "Runs in 2.3 minutes on CPU — no GPU, no API calls", "Performance")
    
    # 3 Stat boxes
    stats = [
        ("28s", "Stage 1: pre-filter\n100K → 3000"),
        ("105s", "Stage 2: embed +\nscore top 3000"),
        ("0s", "Stage 3: reasoning +\nCSV output")
    ]
    
    box_width = Inches(3.5)
    box_height = Inches(3.2)
    left_positions = [Inches(1.0), Inches(4.9), Inches(8.8)]
    
    for idx, (stat, label) in enumerate(stats):
        box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[idx], Inches(1.8), box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = card_bg
        box.line.color.rgb = teal
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p_stat = tf.paragraphs[0]
        p_stat.text = stat
        p_stat.font.name = 'Segoe UI'
        p_stat.font.size = Pt(56)
        p_stat.font.bold = True
        p_stat.font.color.rgb = teal
        p_stat.alignment = PP_ALIGN.CENTER
        p_stat.space_after = Pt(14)
        
        p_lbl = tf.add_paragraph()
        p_lbl.text = label
        p_lbl.font.name = 'Segoe UI'
        p_lbl.font.size = Pt(14)
        p_lbl.font.color.rgb = white
        p_lbl.alignment = PP_ALIGN.CENTER
        
    # Bottom note
    bn = slide7.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.333), Inches(0.8))
    tf_bn = bn.text_frame
    tf_bn.word_wrap = True
    tf_bn.margin_left = tf_bn.margin_top = tf_bn.margin_right = tf_bn.margin_bottom = 0
    p_bn = tf_bn.paragraphs[0]
    p_bn.text = "Well within 5-minute constraint. 16GB RAM. Zero external API calls."
    p_bn.font.name = 'Segoe UI'
    p_bn.font.size = Pt(15)
    p_bn.font.bold = True
    p_bn.font.color.rgb = light_gray
    p_bn.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 8 — Sample Output
    # -------------------------------------------------------------
    slide8 = create_slide()
    add_header(slide8, "What the shortlist looks like", "Sample Output")
    
    table_shape = slide8.shapes.add_table(4, 3, Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.8))
    table = table_shape.table
    table.columns[0].width = Inches(1.0)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(8.833)
    
    headers = ["Rank", "Score", "Why"]
    for i, h_text in enumerate(headers):
        format_cell(table.cell(0, i), h_text, font_size=14, bold=True, text_color=dark_navy, bg_color=teal, align=PP_ALIGN.CENTER if i<2 else PP_ALIGN.LEFT)
        
    data = [
        ("1", "0.94", '"7yr ML Engineer at product startup; shipped embedding-based search to 2M users; open to work, 30d notice, Pune-based"'),
        ("2", "0.91", '"6yr applied ML; vector DB production experience at Series B; strong GitHub activity; willing to relocate to Noida"'),
        ("3", "0.88", '"5yr NLP + retrieval background; evaluation framework experience (NDCG, A/B); concern: 90-day notice period"')
    ]
    
    for row_idx, (rank, score, why) in enumerate(data, 1):
        cell_bg = card_bg if row_idx % 2 == 0 else dark_navy
        format_cell(table.cell(row_idx, 0), rank, font_size=13, bold=True, text_color=white, bg_color=cell_bg, align=PP_ALIGN.CENTER)
        format_cell(table.cell(row_idx, 1), score, font_size=13, bold=True, text_color=teal, bg_color=cell_bg, align=PP_ALIGN.CENTER)
        format_cell(table.cell(row_idx, 2), why, font_size=12, bold=False, text_color=light_gray, bg_color=cell_bg)
        
    callout = slide8.shapes.add_textbox(Inches(1.0), Inches(5.9), Inches(11.333), Inches(0.6))
    tf_c = callout.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    p_c = tf_c.paragraphs[0]
    p_c.text = "Reasoning is specific, JD-connected, and honest about concerns"
    p_c.font.name = 'Segoe UI'
    p_c.font.size = Pt(14)
    p_c.font.bold = True
    p_c.font.color.rgb = teal
    p_c.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 9 — Why This Wins
    # -------------------------------------------------------------
    slide9 = create_slide()
    add_header(slide9, "What separates us from the field", "Why This Wins")
    
    bullets_box = slide9.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    points = [
        ("No API calls", "Runs completely local. Reproducible in sandboxed Docker environment, passes Stage 3 offline checks."),
        ("Honeypot detection", "Timeline consistency checks and expert skill explosion filters eliminate risk of disqualification by impossible profiles."),
        ("Specific reasoning", "Natural language output references actual candidate facts and concerns, passing Stage 4 manual review."),
        ("2.3 min runtime", "Consumes only 46% of the 5-minute hackathon execution budget, leaving massive headroom for future features.")
    ]
    
    for idx, (title, desc) in enumerate(points):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = "✔   " + title + "  —  "
        p.font.name = 'Segoe UI'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = teal
        
        # Append description in white
        run = p.add_run()
        run.text = desc
        run.font.name = 'Segoe UI'
        run.font.size = Pt(15)
        run.font.bold = False
        run.font.color.rgb = white
        
        p.space_after = Pt(28)

    # -------------------------------------------------------------
    # SLIDE 10 — Close
    # -------------------------------------------------------------
    slide10 = create_slide()
    
    # Title & Subtitle centered
    close_box = slide10.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.6))
    tf = close_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "The right engineer is in the dataset."
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = teal
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "We built the system that finds them."
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(22)
    p2.font.color.rgb = white
    p2.alignment = PP_ALIGN.CENTER
    
    # Bottom metadata
    footer_box = slide10.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.333), Inches(0.8))
    tf_f = footer_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
    
    p_f = tf_f.paragraphs[0]
    p_f.text = "GitHub: https://github.com/redrob-hackathon/recruiter-brain"
    p_f.font.name = 'Segoe UI'
    p_f.font.size = Pt(13)
    p_f.font.color.rgb = light_gray
    p_f.alignment = PP_ALIGN.CENTER
    
    p_f2 = tf_f.add_paragraph()
    p_f2.text = "Team Antigravity"
    p_f2.font.name = 'Segoe UI'
    p_f2.font.size = Pt(14)
    p_f2.font.bold = True
    p_f2.font.color.rgb = teal
    p_f2.alignment = PP_ALIGN.CENTER
    p_f2.space_before = Pt(6)

    # Save presentation
    output_path = "redrob_ranker_deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
